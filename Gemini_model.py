# """
# RAG GENERATION & EVALUATION SYSTEM (OPTIMIZED)
# ================================================
# - SentenceTransformer BAAI/bge-m3 on CUDA (shared singleton, never reloaded)
# - Ollama models for LLM generation (gen server: 11435)
# - Gemini 2.0 Flash Lite via GPT-4o-mini endpoint for generation + judge
# - Metrics: BERTScore (GPU), ROUGE, BLEU, METEOR, Faithfulness, Completeness, Correctness
# - Ground truth matching (exact + fuzzy difflib)
# - Multi-model comparison with shared indexing/retrieval

# BUG FIXES:
#   BUG 1: num_ctx=8192 — context no longer silently truncated by Ollama
#   BUG 2: num_predict raised 512→1024; Gemini max_tokens raised
#   BUG 3: upsert() uses PointStruct (not plain dicts)
#   BUG 4: csv, xlsx, json, html, pptx, xml, yaml now supported
#   BUG 5: index_uploaded_files returns chunk count; callers raise on 0
#   BUG 6: __MACOSX and dot-files from Mac ZIPs skipped
#   BUG 7: Context capped at MAX_CONTEXT_CHARS
# """

# import os, json, time, re, math, difflib, requests, logging, threading
# import uuid as _uuid_mod
# from typing import List, Dict, Any, Optional
# from concurrent.futures import ThreadPoolExecutor, as_completed
# from pathlib import Path

# import numpy as np
# from qdrant_client import QdrantClient
# from qdrant_client.models import VectorParams, Distance, PointStruct
# from sentence_transformers import SentenceTransformer
# from sklearn.metrics.pairwise import cosine_similarity

# try:
#     from bert_score import BERTScorer
#     BERTSCORE_AVAILABLE = True
# except ImportError:
#     BERTSCORE_AVAILABLE = False

# try:
#     from rouge_score import rouge_scorer as rouge_scorer_lib
#     ROUGE_AVAILABLE = True
# except ImportError:
#     ROUGE_AVAILABLE = False

# try:
#     from nltk.translate.bleu_score import sentence_bleu, SmoothingFunction
#     from nltk.translate.meteor_score import meteor_score
#     import nltk
#     NLTK_METRICS_AVAILABLE = True
# except ImportError:
#     NLTK_METRICS_AVAILABLE = False

# from Evaluate_Retrieval_With_Reranker import (
#     HybridRetriever, BGE_MODEL, OLLAMA_GEN_URL, OLLAMA_JUDGE_URL, _SHARED_EMBEDDER,
# )

# logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
# logger = logging.getLogger(__name__)

# # ==================================================
# # CONFIG
# # ==================================================
# QDRANT_URL         = "http://localhost:7333"
# TOP_K              = 8
# MAX_CONTEXT_CHARS  = 6000     # BUG 7
# OLLAMA_NUM_CTX     = 8192     # BUG 1
# OLLAMA_NUM_PREDICT = 1024     # BUG 2
# OLLAMA_TIMEOUT     = 300
# JUDGE_TIMEOUT      = 600
# CHUNK_SIZE         = 400
# CHUNK_OVERLAP      = 80

# GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
# if not GEMINI_API_KEY:
#     raise RuntimeError("GEMINI_API_KEY environment variable not set")

# GEMINI_ENDPOINT = (
#     "https://aoai-farm.bosch-temp.com/api/openai/deployments/"
#     "askbosch-prod-farm-openai-gpt-4o-mini-2024-07-18/"
#     "chat/completions?api-version=2024-08-01-preview"
# )

# # ==================================================
# # SINGLETONS
# # ==================================================
# _embedder     = _SHARED_EMBEDDER
# _bert_scorer  = None
# _rouge_scorer = None
# _bert_lock    = threading.Lock()
# _rouge_lock   = threading.Lock()
# print("✓ Embedder ready (reused from retriever module singleton)")


# def _embed(texts: List[str]) -> List[List[float]]:
#     return _embedder.encode(texts, convert_to_numpy=True, show_progress_bar=False, batch_size=64).tolist()


# def get_bert_scorer():
#     global _bert_scorer
#     with _bert_lock:
#         if _bert_scorer is None and BERTSCORE_AVAILABLE:
#             logger.info("Initializing BERTScorer…")
#             _bert_scorer = BERTScorer(lang="en", rescale_with_baseline=True, device="cuda", batch_size=16)
#     return _bert_scorer


# def get_rouge_scorer():
#     global _rouge_scorer
#     with _rouge_lock:
#         if _rouge_scorer is None and ROUGE_AVAILABLE:
#             _rouge_scorer = rouge_scorer_lib.RougeScorer(['rouge1','rouge2','rougeL'], use_stemmer=True)
#     return _rouge_scorer


# # ==================================================
# # JSON UTILS
# # ==================================================

# def clean_json_response(text: str) -> str:
#     if not text or not text.strip():
#         return "{}"
#     text = text.strip()
#     for fence in ("```json", "```"):
#         if text.startswith(fence):
#             text = text[len(fence):]
#     if text.endswith("```"):
#         text = text[:-3]
#     text = text.strip()
#     m = re.search(r'[\[{].*[\]}]', text, re.DOTALL)
#     return m.group(0).strip() if m else text


# def safe_json_parse(text: str, default: Any) -> Any:
#     try:
#         return json.loads(clean_json_response(text))
#     except Exception:
#         return default


# # ==================================================
# # TEXT EXTRACTION
# # ==================================================

# def _extract_text_from_file(file: Path) -> str:
#     if file.name.startswith(".") or file.name.startswith("._") or "__MACOSX" in file.parts:
#         return ""
#     suffix = file.suffix.lower()

#     if suffix in [".txt", ".md", ".yaml", ".yml"]:
#         try: return file.read_text(errors="ignore")
#         except: return ""

#     if suffix == ".pdf":
#         try:
#             import pdfplumber
#             with pdfplumber.open(file) as pdf:
#                 return "\n".join(p.extract_text() for p in pdf.pages if p.extract_text())
#         except Exception as e:
#             logger.warning(f"PDF {file.name}: {e}"); return ""

#     if suffix in [".docx", ".doc"]:
#         try:
#             from docx import Document
#             doc = Document(file); parts = []
#             for p in doc.paragraphs:
#                 if p.text.strip(): parts.append(p.text.strip())
#             for table in doc.tables:
#                 for row in table.rows:
#                     cells = [c.text.strip() for c in row.cells if c.text.strip()]
#                     deduped = []
#                     for c in cells:
#                         if not deduped or c != deduped[-1]: deduped.append(c)
#                     if deduped: parts.append(" | ".join(deduped))
#             return "\n".join(parts)
#         except Exception as e:
#             logger.warning(f"DOCX {file.name}: {e}"); return ""

#     if suffix == ".csv":
#         try:
#             import csv
#             rows = []
#             with open(file, newline="", encoding="utf-8", errors="ignore") as f:
#                 for row in csv.reader(f): rows.append(", ".join(row))
#             return "\n".join(rows)
#         except Exception as e:
#             logger.warning(f"CSV {file.name}: {e}"); return ""

#     if suffix in [".xlsx", ".xls"]:
#         try:
#             import openpyxl
#             wb = openpyxl.load_workbook(file, data_only=True); rows = []
#             for sheet in wb.worksheets:
#                 rows.append(f"[Sheet: {sheet.title}]")
#                 for row in sheet.iter_rows(values_only=True):
#                     cells = [str(c) if c is not None else "" for c in row]
#                     if any(c.strip() for c in cells): rows.append(", ".join(cells))
#             return "\n".join(rows)
#         except Exception as e:
#             logger.warning(f"XLSX {file.name}: {e}"); return ""

#     if suffix == ".json":
#         try:
#             data = json.loads(file.read_text(errors="ignore"))
#             return json.dumps(data, indent=2, ensure_ascii=False)
#         except Exception as e:
#             logger.warning(f"JSON {file.name}: {e}"); return ""

#     if suffix in [".html", ".htm"]:
#         try:
#             from html.parser import HTMLParser
#             class _TE(HTMLParser):
#                 def __init__(self):
#                     super().__init__(); self.parts = []; self._skip = False
#                 def handle_starttag(self, tag, attrs):
#                     if tag in ("script","style"): self._skip = True
#                 def handle_endtag(self, tag):
#                     if tag in ("script","style"): self._skip = False
#                 def handle_data(self, data):
#                     if not self._skip and data.strip(): self.parts.append(data.strip())
#             p = _TE(); p.feed(file.read_text(errors="ignore"))
#             return " ".join(p.parts)
#         except Exception as e:
#             logger.warning(f"HTML {file.name}: {e}"); return ""

#     if suffix in [".pptx", ".ppt"]:
#         try:
#             from pptx import Presentation
#             prs = Presentation(file); parts = []
#             for i, slide in enumerate(prs.slides, 1):
#                 parts.append(f"[Slide {i}]")
#                 for shape in slide.shapes:
#                     if hasattr(shape,"text") and shape.text.strip(): parts.append(shape.text.strip())
#             return "\n".join(parts)
#         except Exception as e:
#             logger.warning(f"PPTX {file.name}: {e}"); return ""

#     if suffix == ".xml":
#         try:
#             import xml.etree.ElementTree as ET
#             root = ET.parse(file).getroot()
#             return "\n".join(el.text.strip() for el in root.iter() if el.text and el.text.strip())
#         except Exception as e:
#             logger.warning(f"XML {file.name}: {e}"); return ""

#     return ""


# # ==================================================
# # COLLECTION HELPERS
# # ==================================================

# def create_temp_collection(collection_name: str) -> QdrantClient:
#     client = QdrantClient(url=QDRANT_URL)
#     try:
#         if client.collection_exists(collection_name): client.delete_collection(collection_name)
#     except: pass
#     client.create_collection(collection_name=collection_name,
#                              vectors_config=VectorParams(size=1024, distance=Distance.COSINE))
#     return client


# def chunk_text(text: str, filename: str) -> List[Dict]:
#     words = text.split(); chunks = []; start = 0
#     while start < len(words):
#         end = min(start + CHUNK_SIZE, len(words))
#         chunks.append({"text": " ".join(words[start:end]), "filename": filename,
#                         "chunk_idx": len(chunks), "word_start": start, "word_end": end})
#         if end == len(words): break
#         start += CHUNK_SIZE - CHUNK_OVERLAP
#     return chunks


# def index_uploaded_files(client: QdrantClient, collection_name: str, docs_path) -> int:
#     """Returns number of indexed chunks (BUG 3+5+6 fix)."""
#     all_chunks: List[Dict] = []
#     for file in Path(docs_path).rglob("*"):
#         if not file.is_file(): continue
#         if any(p.startswith(".") or p == "__MACOSX" for p in file.parts): continue
#         text = _extract_text_from_file(file)
#         if not text.strip(): continue
#         file_chunks = chunk_text(text, file.name)
#         all_chunks.extend(file_chunks)
#         logger.info(f"  {file.name}: {len(file_chunks)} chunks")

#     if not all_chunks:
#         logger.warning("No indexable content found."); return 0

#     logger.info(f"Encoding {len(all_chunks)} chunks…")
#     embeddings = _embed([c["text"] for c in all_chunks])

#     # BUG 3 FIX: PointStruct
#     points = [
#         PointStruct(id=i, vector=emb,
#                     payload={"content": c["text"], "filename": c["filename"], "chunk_idx": c["chunk_idx"]})
#         for i, (c, emb) in enumerate(zip(all_chunks, embeddings))
#     ]
#     for i in range(0, len(points), 100):
#         client.upsert(collection_name=collection_name, points=points[i:i+100])

#     logger.info(f"✓ Indexed {len(points)} chunks from {len(set(c['filename'] for c in all_chunks))} files")
#     return len(points)


# def delete_temp_collection(collection_name: str) -> None:
#     try: QdrantClient(url=QDRANT_URL).delete_collection(collection_name)
#     except: pass


# # ==================================================
# # GENERATION
# # ==================================================

# def get_available_ollama_models() -> List[str]:
#     try:
#         r = requests.get(f"{OLLAMA_GEN_URL}/api/tags", timeout=5)
#         if r.status_code == 200:
#             return [m["name"] for m in r.json().get("models", [])]
#     except Exception as e:
#         logger.warning(f"Ollama unavailable: {e}")
#     return []


# def generate_with_ollama(model_name: str, prompt: str, temperature: float = 0.1) -> str:
#     """BUG 1+2 FIX: num_ctx=8192, num_predict=1024."""
#     payload = {"model": model_name, "prompt": prompt, "temperature": temperature, "stream": False,
#                "options": {"num_ctx": OLLAMA_NUM_CTX, "num_predict": OLLAMA_NUM_PREDICT, "top_p": 0.9}}
#     try:
#         r = requests.post(f"{OLLAMA_GEN_URL}/api/generate", json=payload, timeout=OLLAMA_TIMEOUT)
#         return r.json()["response"] if r.status_code == 200 else f"Error: {r.status_code}"
#     except Exception as e:
#         return f"Error: {e}"


# def generate_with_gemini(prompt: str, temperature: float = 0.1) -> str:
#     headers = {"genaiplatform-farm-subscription-key": GEMINI_API_KEY, "Content-Type": "application/json"}
#     payload = {"messages": [{"role":"user","content":prompt}], "temperature": temperature, "max_tokens": 1024}
#     try:
#         r = requests.post(GEMINI_ENDPOINT, headers=headers, json=payload, timeout=60)
#         return r.json()["choices"][0]["message"]["content"] if r.status_code == 200 else f"Error: {r.status_code}"
#     except Exception as e:
#         return f"Error: {e}"


# # ==================================================
# # JUDGE
# # ==================================================

# def evaluate_with_gpt4o_mini(question: str, context: str, answer: str, ground_truth: str) -> Dict[str, Any]:
#     has_gt = bool(ground_truth.strip())
#     gt_section = (f"**Ground Truth Answer:** {ground_truth}"
#                   if has_gt else "**Ground Truth:** Not available — evaluate against context only.")
#     judge_prompt = f"""You are a QA Auditor for a RAG system.

# **Question:** {question}
# **Retrieved Context:** {context[:1500]}
# **Generated Answer:** {answer}
# {gt_section}

# Score 0-10 integers:
# - faithfulness: Are all claims supported by context?
# - completeness: Are all aspects of the question addressed?
# - correctness: {"Accuracy vs ground truth?" if has_gt else "Factual reasonableness given context?"}

# Return ONLY valid JSON (no markdown):
# {{"faithfulness":<int>,"completeness":<int>,"correctness":<int>,"explanation":"<one sentence>"}}"""

#     headers = {"genaiplatform-farm-subscription-key": GEMINI_API_KEY, "Content-Type": "application/json"}
#     payload = {"messages": [{"role":"user","content":judge_prompt}], "temperature": 0.1, "max_tokens": 300}
#     default = {"faithfulness": 0, "completeness": 0, "correctness": 0, "explanation": "Judge failed"}
#     try:
#         r = requests.post(GEMINI_ENDPOINT, headers=headers, json=payload, timeout=60)
#         if r.status_code == 200:
#             scores = safe_json_parse(r.json()["choices"][0]["message"]["content"], default)
#             for k in ["faithfulness","completeness","correctness"]:
#                 scores[k] = max(0, min(10, int(scores.get(k, 0))))
#             return scores
#     except Exception as e:
#         logger.error(f"Judge error: {e}")
#     return default


# # ==================================================
# # METRICS
# # ==================================================

# def calculate_fast_metrics(generated: str, ground_truth: str) -> Dict[str, float]:
#     metrics: Dict[str, float] = {}

#     if BERTSCORE_AVAILABLE:
#         try:
#             P, R, F1 = get_bert_scorer().score([generated], [ground_truth])
#             metrics.update({"bertscore_precision": float(P[0]), "bertscore_recall": float(R[0]),
#                              "bertscore_f1": float(F1[0])})
#         except Exception as e:
#             logger.warning(f"BERTScore error: {e}")
#             metrics.update({"bertscore_precision": 0.0, "bertscore_recall": 0.0, "bertscore_f1": 0.0})
#     else:
#         metrics.update({"bertscore_precision": 0.0, "bertscore_recall": 0.0, "bertscore_f1": 0.0})

#     if ROUGE_AVAILABLE:
#         try:
#             rouge = get_rouge_scorer().score(ground_truth, generated)
#             metrics.update({"rouge1_f1": rouge['rouge1'].fmeasure, "rouge2_f1": rouge['rouge2'].fmeasure,
#                              "rougeL_f1": rouge['rougeL'].fmeasure})
#         except Exception as e:
#             logger.warning(f"ROUGE error: {e}")
#             metrics.update({"rouge1_f1": 0.0, "rouge2_f1": 0.0, "rougeL_f1": 0.0})
#     else:
#         metrics.update({"rouge1_f1": 0.0, "rouge2_f1": 0.0, "rougeL_f1": 0.0})

#     # BLEU
#     if NLTK_METRICS_AVAILABLE and generated.strip() and ground_truth.strip():
#         try:
#             metrics["bleu_score"] = float(sentence_bleu(
#                 [ground_truth.split()], generated.split(), smoothing_function=SmoothingFunction().method4))
#         except: metrics["bleu_score"] = 0.0
#     else:
#         metrics["bleu_score"] = 0.0

#     # METEOR
#     if NLTK_METRICS_AVAILABLE and generated.strip() and ground_truth.strip():
#         try: metrics["meteor_score"] = float(meteor_score([ground_truth.split()], generated.split()))
#         except: metrics["meteor_score"] = 0.0
#     else:
#         metrics["meteor_score"] = 0.0

#     return metrics


# def calculate_semantic_similarity(a: str, b: str) -> float:
#     if not a.strip() or not b.strip(): return 0.0
#     try:
#         e1 = np.array(_embed([a])[0]).reshape(1,-1)
#         e2 = np.array(_embed([b])[0]).reshape(1,-1)
#         return float(cosine_similarity(e1, e2)[0][0])
#     except: return 0.0


# # ==================================================
# # RAG SYSTEM
# # ==================================================

# class RAGSystem:
#     def __init__(self, qdrant_url: str, collection: str):
#         self.retriever = HybridRetriever(
#             qdrant_url=qdrant_url, collection_name=collection,
#             use_ollama=True, use_reranker=True, embedder=_embedder,
#         )

#     def retrieve(self, query: str, top_k: int = TOP_K):
#         results = self.retriever.search(query, top_k=top_k)
#         docs = [{"content": r.content, "score": r.score,
#                  "source": r.metadata.get("filename","unknown"),
#                  "chunk":  r.metadata.get("chunk_idx", 0)} for r in results]
#         sources = [d["source"] for d in docs]
#         parts   = [f"[Source {i}: {d['source']} | chunk {d['chunk']} | score {d['score']:.4f}]\n{d['content']}"
#                    for i, d in enumerate(docs, 1)]
#         context = "\n\n---\n\n".join(parts)
#         if len(context) > MAX_CONTEXT_CHARS:       # BUG 7 FIX
#             context = context[:MAX_CONTEXT_CHARS] + "\n\n[Context truncated]"
#         return docs, context, sources

#     def build_prompt(self, question: str, context: str) -> str:
#         return (
#             "You are a knowledgeable assistant. Answer using the context below.\n\n"
#             "IMPORTANT: Read ALL chunks carefully. The answer IS in the context.\n"
#             "Give a direct, complete answer. Only say 'not found' if truly absent.\n\n"
#             f"CONTEXT:\n{context}\n\nQUESTION: {question}\n\nAnswer:"
#         )


# # ==================================================
# # GROUND TRUTH MATCHING
# # ==================================================

# def find_ground_truth(question: str, gt_data: dict) -> Optional[dict]:
#     """Exact → case-insensitive → fuzzy (threshold 0.75)."""
#     if not gt_data or not gt_data.get("questions"):
#         return None
#     qs = gt_data["questions"]
#     qn = question.strip().lower()

#     for e in qs:
#         if e.get("question","").strip() == question.strip():
#             return e
#     for e in qs:
#         if e.get("question","").strip().lower() == qn:
#             return e

#     best, best_e = 0.0, None
#     for e in qs:
#         r = difflib.SequenceMatcher(None, qn, e.get("question","").lower()).ratio()
#         if r > best: best, best_e = r, e

#     if best >= 0.75:
#         logger.info(f"GT fuzzy match ratio={best:.3f}")
#         return best_e
#     return None


# # ==================================================
# # SHARED EVALUATION HELPER
# # ==================================================

# def _build_eval_result(
#     model_key: str, answer: str, latency: float,
#     question: str, context: str, reference: str,
#     docs: list, retrieved_sources: list,
# ) -> dict:
#     """Evaluate answer against reference using judge + BERTScore + semantic sim."""
#     def _judge():
#         return evaluate_with_gpt4o_mini(question=question, context=context,
#                                         answer=answer, ground_truth=reference)
#     def _fast():
#         return calculate_fast_metrics(answer, reference) if reference.strip() else {}
#     def _sem():
#         return calculate_semantic_similarity(answer, reference) if reference.strip() else 0.0

#     with ThreadPoolExecutor(max_workers=3) as ex:
#         fj = ex.submit(_judge); ff = ex.submit(_fast); fs = ex.submit(_sem)
#         judge = fj.result(); fast = ff.result(); sem = fs.result()

#     overall = judge.get("faithfulness",0)*0.35 + judge.get("completeness",0)*0.35 + judge.get("correctness",0)*0.30
#     uniq    = list(dict.fromkeys(retrieved_sources))

#     return {
#         "model":               model_key,
#         "answer":              answer,
#         "latency_sec":         round(latency, 2),
#         "faithfulness":        judge.get("faithfulness", 0),
#         "completeness":        judge.get("completeness", 0),
#         "correctness":         judge.get("correctness",  0),
#         "overall_score":       round(overall, 2),
#         "judge_explanation":   judge.get("explanation",""),
#         "semantic_similarity": round(sem, 4),
#         "bertscore_f1":        round(fast.get("bertscore_f1",        0.0), 4),
#         "bertscore_precision": round(fast.get("bertscore_precision",  0.0), 4),
#         "bertscore_recall":    round(fast.get("bertscore_recall",     0.0), 4),
#         "rouge1_f1":           round(fast.get("rouge1_f1",  0.0), 4),
#         "rouge2_f1":           round(fast.get("rouge2_f1",  0.0), 4),
#         "rougeL_f1":           round(fast.get("rougeL_f1",  0.0), 4),
#         "bleu_score":          round(fast.get("bleu_score", 0.0), 4),
#         "meteor_score":        round(fast.get("meteor_score",0.0), 4),
#         "chunks_retrieved":    len(retrieved_sources),
#         "unique_files":        len(uniq),
#         "avg_chunk_score":     round(sum(d["score"] for d in docs)/len(docs), 4) if docs else 0.0,
#         "top_file":            uniq[0] if uniq else "none",
#     }


# # ==================================================
# # SINGLE-MODEL COMPARE (existing app flow)
# # ==================================================

# def run_rag_evaluation(
#     selected_model: str, user_question: str, collection_name: str,
#     docs_path, gt_entry: Optional[dict] = None,
# ) -> Dict[str, Any]:
#     client      = create_temp_collection(collection_name)
#     chunk_count = index_uploaded_files(client, collection_name, docs_path)
#     if chunk_count == 0:
#         raise ValueError("No content extracted. Supported: PDF, DOCX, TXT, MD, CSV, XLSX, JSON, HTML, PPTX, XML, YAML.")

#     rag  = RAGSystem(QDRANT_URL, collection_name)
#     docs, context, sources = rag.retrieve(user_question, top_k=TOP_K)
#     if not docs:
#         raise ValueError("Retrieval returned no results.")

#     prompt = rag.build_prompt(user_question, context)

#     model_answer = ""; gemini_answer = ""; ml = 0.0; gl = 0.0

#     def _om():
#         nonlocal model_answer, ml
#         t0 = time.time(); model_answer  = generate_with_ollama(selected_model, prompt); ml = time.time()-t0
#     def _gm():
#         nonlocal gemini_answer, gl
#         t0 = time.time(); gemini_answer = generate_with_gemini(prompt);                gl = time.time()-t0

#     with ThreadPoolExecutor(max_workers=2) as ex:
#         f1 = ex.submit(_om); f2 = ex.submit(_gm); f1.result(); f2.result()

#     gt_ans      = gt_entry.get("answer","") if gt_entry else ""
#     reference   = gt_ans if gt_ans else gemini_answer
#     ref_label   = "Ground Truth" if gt_ans else "Gemini Answer"
#     eval_answer = model_answer if model_answer else gemini_answer

#     res = _build_eval_result(selected_model, eval_answer, ml, user_question, context, reference, docs, sources)

#     # Semantic similarity vs Gemini (always)
#     sem_vs_gemini = calculate_semantic_similarity(eval_answer, gemini_answer) if model_answer else 1.0

#     uniq = list(dict.fromkeys(sources))
#     delete_temp_collection(collection_name)

#     metrics = {**res,
#                "semantic_similarity_vs_gemini": round(sem_vs_gemini, 4),
#                "reference_label":  ref_label,
#                "model_latency_sec": round(ml, 2),
#                "gemini_latency_sec": round(gl, 2),
#                "has_ground_truth": bool(gt_ans)}
#     metrics.pop("answer", None)

#     return {"model_answer": model_answer, "gemini_answer": gemini_answer,
#             "model_sources": sources, "metrics": metrics, "gt_entry": gt_entry}


# # ==================================================
# # MULTI-MODEL COMPARISON
# # ==================================================

# def run_multi_model_comparison(
#     model_names: List[str], user_question: str, docs_path,
#     gt_entry: Optional[dict] = None,
# ) -> Dict[str, Any]:
#     """
#     Index once, retrieve once, generate all models (Ollama sequential + Gemini parallel),
#     evaluate all in parallel. Returns per-model results dict.
#     """
#     gt_answer = gt_entry.get("answer","") if gt_entry else ""
#     has_gt    = bool(gt_answer.strip())

#     # Index
#     cname  = f"temp_multi_{_uuid_mod.uuid4().hex[:8]}"
#     client = create_temp_collection(cname)
#     n      = index_uploaded_files(client, cname, docs_path)
#     if n == 0:
#         delete_temp_collection(cname)
#         raise ValueError("No content extracted from uploaded files.")

#     # Retrieve
#     rag = RAGSystem(QDRANT_URL, cname)
#     docs, context, sources = rag.retrieve(user_question, top_k=TOP_K)
#     if not docs:
#         delete_temp_collection(cname)
#         raise ValueError("Retrieval returned no results.")

#     prompt = rag.build_prompt(user_question, context)

#     # Generate — Gemini async, Ollama sequential
#     generation: Dict[str, Dict] = {}

#     def _gen_gemini():
#         t0 = time.time(); a = generate_with_gemini(prompt)
#         return {"answer": a, "latency": time.time()-t0}

#     with ThreadPoolExecutor(max_workers=1) as gex:
#         gf = gex.submit(_gen_gemini)
#         for name in model_names:
#             logger.info(f"Generating: {name}")
#             try:
#                 t0 = time.time(); ans = generate_with_ollama(name, prompt)
#                 generation[name] = {"answer": ans, "latency": time.time()-t0}
#             except Exception as e:
#                 generation[name] = {"answer": f"Error: {e}", "latency": 0.0}
#         try:
#             generation["Gemini"] = gf.result()
#         except Exception as e:
#             generation["Gemini"] = {"answer": f"Error: {e}", "latency": 0.0}

#     gemini_answer = generation["Gemini"]["answer"]
#     ref_ollama    = gt_answer if has_gt else gemini_answer   # for Ollama models
#     ref_gemini    = gt_answer                                  # for Gemini (empty = context-only judge)

#     # Evaluate all in parallel
#     eval_results: Dict[str, Dict] = {}

#     def _eval(key: str):
#         gen = generation.get(key, {})
#         ref = ref_gemini if key == "Gemini" else ref_ollama
#         return _build_eval_result(key, gen.get("answer",""), gen.get("latency",0.0),
#                                    user_question, context, ref, docs, sources)

#     all_keys = list(model_names) + ["Gemini"]
#     with ThreadPoolExecutor(max_workers=min(len(all_keys), 3)) as ex:
#         futures = {ex.submit(_eval, k): k for k in all_keys}
#         for f in as_completed(futures):
#             k = futures[f]
#             try:    eval_results[k] = f.result()
#             except Exception as e:
#                 logger.error(f"Eval failed {k}: {e}")
#                 eval_results[k] = {"model": k, "overall_score": 0, "error": str(e)}

#     delete_temp_collection(cname)

#     return {
#         "results":           eval_results,
#         "gemini_answer":     gemini_answer,
#         "ground_truth":      gt_entry,
#         "has_ground_truth":  has_gt,
#         "reference_label":   "Ground Truth" if has_gt else "Gemini Answer",
#         "retrieved_sources": sources,
#     }



"""
RAG GENERATION & EVALUATION SYSTEM (OPTIMIZED)
================================================
- SentenceTransformer BAAI/bge-m3 on CUDA (shared singleton, never reloaded)
- Ollama models for LLM generation (gen server: 11435)
- Gemini 2.0 Flash Lite via GPT-4o-mini endpoint for generation + judge
- Metrics: BERTScore (GPU), ROUGE, BLEU, METEOR, Faithfulness, Completeness, Correctness
- Ground truth matching (exact + fuzzy difflib)
- Multi-model comparison with shared indexing/retrieval

BUG FIXES:
  BUG 1: num_ctx=8192 — context no longer silently truncated by Ollama
  BUG 2: num_predict raised 512→1024; Gemini max_tokens raised
  BUG 3: upsert() uses PointStruct (not plain dicts)
  BUG 4: csv, xlsx, json, html, pptx, xml, yaml now supported
  BUG 5: index_uploaded_files returns chunk count; callers raise on 0
  BUG 6: __MACOSX and dot-files from Mac ZIPs skipped
  BUG 7: Context capped at MAX_CONTEXT_CHARS
"""

import os, json, time, re, math, difflib, requests, logging, threading
import uuid as _uuid_mod
from typing import List, Dict, Any, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

import numpy as np
from qdrant_client import QdrantClient
from qdrant_client.models import VectorParams, Distance, PointStruct
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity

try:
    from bert_score import BERTScorer
    BERTSCORE_AVAILABLE = True
except ImportError:
    BERTSCORE_AVAILABLE = False

try:
    from rouge_score import rouge_scorer as rouge_scorer_lib
    ROUGE_AVAILABLE = True
except ImportError:
    ROUGE_AVAILABLE = False

try:
    from nltk.translate.bleu_score import sentence_bleu, SmoothingFunction
    from nltk.translate.meteor_score import meteor_score
    import nltk
    NLTK_METRICS_AVAILABLE = True
except ImportError:
    NLTK_METRICS_AVAILABLE = False

from Evaluate_Retrieval_With_Reranker import (
    HybridRetriever, BGE_MODEL, OLLAMA_GEN_URL, OLLAMA_JUDGE_URL, _SHARED_EMBEDDER,
)

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ==================================================
# CONFIG
# ==================================================
QDRANT_URL         = "http://localhost:7333"
TOP_K              = 8
MAX_CONTEXT_CHARS  = 14000        # BUG 7
OLLAMA_NUM_CTX     = 8192      # BUG 1
OLLAMA_NUM_PREDICT = 2048     # BUG 2
OLLAMA_TIMEOUT     = 300
JUDGE_TIMEOUT      = 600
CHUNK_SIZE         = 300   
CHUNK_OVERLAP      = 100   

GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    raise RuntimeError("GEMINI_API_KEY environment variable not set")

GEMINI_ENDPOINT = (
    "https://aoai-farm.bosch-temp.com/api/openai/deployments/"
    "askbosch-prod-farm-openai-gpt-4o-mini-2024-07-18/"
    "chat/completions?api-version=2024-08-01-preview"
)

# ==================================================
# SINGLETONS
# ==================================================
_embedder     = _SHARED_EMBEDDER
_bert_scorer  = None
_rouge_scorer = None
_bert_lock    = threading.Lock()
_rouge_lock   = threading.Lock()
print("✓ Embedder ready (reused from retriever module singleton)")


def _embed(texts: List[str]) -> List[List[float]]:
    return _embedder.encode(texts, convert_to_numpy=True, show_progress_bar=False, batch_size=64).tolist()


def get_bert_scorer():
    global _bert_scorer
    with _bert_lock:
        if _bert_scorer is None and BERTSCORE_AVAILABLE:
            logger.info("Initializing BERTScorer…")
            _bert_scorer = BERTScorer(lang="en", rescale_with_baseline=True, device="cuda", batch_size=16)
    return _bert_scorer


def get_rouge_scorer():
    global _rouge_scorer
    with _rouge_lock:
        if _rouge_scorer is None and ROUGE_AVAILABLE:
            _rouge_scorer = rouge_scorer_lib.RougeScorer(['rouge1','rouge2','rougeL'], use_stemmer=True)
    return _rouge_scorer


# ==================================================
# JSON UTILS
# ==================================================

def clean_json_response(text: str) -> str:
    if not text or not text.strip():
        return "{}"
    text = text.strip()
    for fence in ("```json", "```"):
        if text.startswith(fence):
            text = text[len(fence):]
    if text.endswith("```"):
        text = text[:-3]
    text = text.strip()
    m = re.search(r'[\[{].*[\]}]', text, re.DOTALL)
    return m.group(0).strip() if m else text


def safe_json_parse(text: str, default: Any) -> Any:
    try:
        return json.loads(clean_json_response(text))
    except Exception:
        return default


# ==================================================
# TEXT EXTRACTION
# ==================================================

def _extract_text_from_file(file: Path) -> str:
    if file.name.startswith(".") or file.name.startswith("._") or "__MACOSX" in file.parts:
        return ""
    suffix = file.suffix.lower()

    if suffix in [".txt", ".md", ".yaml", ".yml"]:
        try: return file.read_text(errors="ignore")
        except: return ""

    if suffix == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(file) as pdf:
                return "\n".join(p.extract_text() for p in pdf.pages if p.extract_text())
        except Exception as e:
            logger.warning(f"PDF {file.name}: {e}"); return ""

    if suffix in [".docx", ".doc"]:
        try:
            from docx import Document
            doc = Document(file); parts = []
            for p in doc.paragraphs:
                if p.text.strip(): parts.append(p.text.strip())
            for table in doc.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    deduped = []
                    for c in cells:
                        if not deduped or c != deduped[-1]: deduped.append(c)
                    if deduped: parts.append(" | ".join(deduped))
            return "\n".join(parts)
        except Exception as e:
            logger.warning(f"DOCX {file.name}: {e}"); return ""

    if suffix == ".csv":
        try:
            import csv
            rows = []
            with open(file, newline="", encoding="utf-8", errors="ignore") as f:
                for row in csv.reader(f): rows.append(", ".join(row))
            return "\n".join(rows)
        except Exception as e:
            logger.warning(f"CSV {file.name}: {e}"); return ""

    if suffix in [".xlsx", ".xls"]:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file, data_only=True); rows = []
            for sheet in wb.worksheets:
                rows.append(f"[Sheet: {sheet.title}]")
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(c.strip() for c in cells): rows.append(", ".join(cells))
            return "\n".join(rows)
        except Exception as e:
            logger.warning(f"XLSX {file.name}: {e}"); return ""

    if suffix == ".json":
        try:
            data = json.loads(file.read_text(errors="ignore"))
            return json.dumps(data, indent=2, ensure_ascii=False)
        except Exception as e:
            logger.warning(f"JSON {file.name}: {e}"); return ""

    if suffix in [".html", ".htm"]:
        try:
            from html.parser import HTMLParser
            class _TE(HTMLParser):
                def __init__(self):
                    super().__init__(); self.parts = []; self._skip = False
                def handle_starttag(self, tag, attrs):
                    if tag in ("script","style"): self._skip = True
                def handle_endtag(self, tag):
                    if tag in ("script","style"): self._skip = False
                def handle_data(self, data):
                    if not self._skip and data.strip(): self.parts.append(data.strip())
            p = _TE(); p.feed(file.read_text(errors="ignore"))
            return " ".join(p.parts)
        except Exception as e:
            logger.warning(f"HTML {file.name}: {e}"); return ""

    if suffix in [".pptx", ".ppt"]:
        try:
            from pptx import Presentation
            prs = Presentation(file); parts = []
            for i, slide in enumerate(prs.slides, 1):
                parts.append(f"[Slide {i}]")
                for shape in slide.shapes:
                    if hasattr(shape,"text") and shape.text.strip(): parts.append(shape.text.strip())
            return "\n".join(parts)
        except Exception as e:
            logger.warning(f"PPTX {file.name}: {e}"); return ""

    if suffix == ".xml":
        try:
            import xml.etree.ElementTree as ET
            root = ET.parse(file).getroot()
            return "\n".join(el.text.strip() for el in root.iter() if el.text and el.text.strip())
        except Exception as e:
            logger.warning(f"XML {file.name}: {e}"); return ""

    return ""


# ==================================================
# COLLECTION HELPERS
# ==================================================

def create_temp_collection(collection_name: str) -> QdrantClient:
    client = QdrantClient(url=QDRANT_URL)
    try:
        if client.collection_exists(collection_name): client.delete_collection(collection_name)
    except: pass
    client.create_collection(collection_name=collection_name,
                             vectors_config=VectorParams(size=1024, distance=Distance.COSINE))
    return client


def chunk_text(text: str, filename: str) -> List[Dict]:
    words = text.split(); chunks = []; start = 0
    while start < len(words):
        end = min(start + CHUNK_SIZE, len(words))
        chunks.append({"text": " ".join(words[start:end]), "filename": filename,
                        "chunk_idx": len(chunks), "word_start": start, "word_end": end})
        if end == len(words): break
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return chunks


def index_uploaded_files(client: QdrantClient, collection_name: str, docs_path) -> int:
    """Returns number of indexed chunks (BUG 3+5+6 fix)."""
    all_chunks: List[Dict] = []
    for file in Path(docs_path).rglob("*"):
        if not file.is_file(): continue
        if any(p.startswith(".") or p == "__MACOSX" for p in file.parts): continue
        text = _extract_text_from_file(file)
        if not text.strip(): continue
        file_chunks = chunk_text(text, file.name)
        all_chunks.extend(file_chunks)
        logger.info(f"  {file.name}: {len(file_chunks)} chunks")

    if not all_chunks:
        logger.warning("No indexable content found."); return 0

    logger.info(f"Encoding {len(all_chunks)} chunks…")
    embeddings = _embed([c["text"] for c in all_chunks])

    # BUG 3 FIX: PointStruct
    points = [
        PointStruct(id=i, vector=emb,
                    payload={"content": c["text"], "filename": c["filename"], "chunk_idx": c["chunk_idx"]})
        for i, (c, emb) in enumerate(zip(all_chunks, embeddings))
    ]
    for i in range(0, len(points), 100):
        client.upsert(collection_name=collection_name, points=points[i:i+100])

    logger.info(f"✓ Indexed {len(points)} chunks from {len(set(c['filename'] for c in all_chunks))} files")
    return len(points)


def delete_temp_collection(collection_name: str) -> None:
    try: QdrantClient(url=QDRANT_URL).delete_collection(collection_name)
    except: pass


# ==================================================
# GENERATION
# ==================================================

def get_available_ollama_models() -> List[str]:
    try:
        r = requests.get(f"{OLLAMA_GEN_URL}/api/tags", timeout=5)
        if r.status_code == 200:
            return [m["name"] for m in r.json().get("models", [])]
    except Exception as e:
        logger.warning(f"Ollama unavailable: {e}")
    return []


def generate_with_ollama(model_name: str, prompt: str, temperature: float = 0.1) -> str:
    """BUG 1+2 FIX: num_ctx=8192, num_predict=1024."""
    payload = {"model": model_name, "prompt": prompt, "temperature": temperature, "stream": False,
               "options": {"num_ctx": OLLAMA_NUM_CTX, "num_predict": OLLAMA_NUM_PREDICT, "top_p": 0.9}}
    try:
        r = requests.post(f"{OLLAMA_GEN_URL}/api/generate", json=payload, timeout=OLLAMA_TIMEOUT)
        return r.json()["response"] if r.status_code == 200 else f"Error: {r.status_code}"
    except Exception as e:
        return f"Error: {e}"


def generate_with_gemini(prompt: str, temperature: float = 0.1) -> str:
    headers = {"genaiplatform-farm-subscription-key": GEMINI_API_KEY, "Content-Type": "application/json"}
    payload = {"messages": [{"role":"user","content":prompt}], "temperature": temperature, "max_tokens": 2048}
    try:
        r = requests.post(GEMINI_ENDPOINT, headers=headers, json=payload, timeout=60)
        return r.json()["choices"][0]["message"]["content"] if r.status_code == 200 else f"Error: {r.status_code}"
    except Exception as e:
        return f"Error: {e}"


# ==================================================
# JUDGE
# ==================================================

def evaluate_with_gpt4o_mini(question: str, context: str, answer: str, ground_truth: str) -> Dict[str, Any]:
    has_gt = bool(ground_truth.strip())
    gt_section = (f"**Ground Truth Answer:** {ground_truth}"
                  if has_gt else "**Ground Truth:** Not available — evaluate against context only.")
    judge_prompt = f"""You are a QA Auditor for a RAG system.

**Question:** {question}
**Retrieved Context:** {context[:1500]}
**Generated Answer:** {answer}
{gt_section}

Score 0-10 integers:
- faithfulness: Are all claims supported by context?
- completeness: Are all aspects of the question addressed?
- correctness: {"Accuracy vs ground truth?" if has_gt else "Factual reasonableness given context?"}

Return ONLY valid JSON (no markdown):
{{"faithfulness":<int>,"completeness":<int>,"correctness":<int>,"explanation":"<one sentence>"}}"""

    headers = {"genaiplatform-farm-subscription-key": GEMINI_API_KEY, "Content-Type": "application/json"}
    payload = {"messages": [{"role":"user","content":judge_prompt}], "temperature": 0.1, "max_tokens": 300}
    default = {"faithfulness": 0, "completeness": 0, "correctness": 0, "explanation": "Judge failed"}
    try:
        r = requests.post(GEMINI_ENDPOINT, headers=headers, json=payload, timeout=60)
        if r.status_code == 200:
            scores = safe_json_parse(r.json()["choices"][0]["message"]["content"], default)
            for k in ["faithfulness","completeness","correctness"]:
                scores[k] = max(0, min(10, int(scores.get(k, 0))))
            return scores
    except Exception as e:
        logger.error(f"Judge error: {e}")
    return default


# ==================================================
# JUDGE ROUTING — select GPT or any Ollama model as judge
# ==================================================

JUDGE_GPT = "GPT-4o-mini"   # sentinel value for the Bosch/GPT endpoint


def evaluate_with_ollama_judge_model(
    model_name: str, question: str, context: str, answer: str, ground_truth: str
) -> Dict[str, Any]:
    """Run the judge prompt on a chosen Ollama model instead of GPT-4o-mini."""
    has_gt = bool(ground_truth.strip())
    gt_section = (f"**Ground Truth Answer:** {ground_truth}"
                  if has_gt else "**Ground Truth:** Not available — evaluate against context only.")
    judge_prompt = f"""You are a QA Auditor for a RAG system.

**Question:** {question}
**Retrieved Context:** {context[:1500]}
**Generated Answer:** {answer}
{gt_section}

Score 0-10 integers:
- faithfulness: Are all claims supported by context?
- completeness: Are all aspects of the question addressed?
- correctness: {"Accuracy vs ground truth?" if has_gt else "Factual reasonableness given context?"}

Return ONLY valid JSON (no markdown, no extra text):
{{"faithfulness":<int>,"completeness":<int>,"correctness":<int>,"explanation":"<one sentence>"}}"""

    payload = {
        "model": model_name, "prompt": judge_prompt, "temperature": 0.0, "stream": False,
        "options": {"num_ctx": 4096, "num_predict": 300},
    }
    default = {"faithfulness": 0, "completeness": 0, "correctness": 0, "explanation": "Judge failed"}
    try:
        r = requests.post(f"{OLLAMA_JUDGE_URL}/api/generate", json=payload, timeout=JUDGE_TIMEOUT)
        if r.status_code == 200:
            scores = safe_json_parse(r.json()["response"], default)
            for k in ["faithfulness", "completeness", "correctness"]:
                scores[k] = max(0, min(10, int(scores.get(k, 0))))
            if "explanation" not in scores:
                scores["explanation"] = ""
            return scores
    except Exception as e:
        logger.error(f"Ollama judge ({model_name}) error: {e}")
    return default


def evaluate_with_judge(
    judge_model: str, question: str, context: str, answer: str, ground_truth: str
) -> Dict[str, Any]:
    """Dispatcher: routes to GPT-4o-mini or the chosen Ollama judge model."""
    if not judge_model or judge_model == JUDGE_GPT:
        return evaluate_with_gpt4o_mini(question, context, answer, ground_truth)
    return evaluate_with_ollama_judge_model(judge_model, question, context, answer, ground_truth)


# ==================================================
# METRICS
# ==================================================

def calculate_fast_metrics(generated: str, ground_truth: str) -> Dict[str, float]:
    metrics: Dict[str, float] = {}

    if BERTSCORE_AVAILABLE:
        try:
            P, R, F1 = get_bert_scorer().score([generated], [ground_truth])
            metrics.update({"bertscore_precision": float(P[0]), "bertscore_recall": float(R[0]),
                             "bertscore_f1": float(F1[0])})
        except Exception as e:
            logger.warning(f"BERTScore error: {e}")
            metrics.update({"bertscore_precision": 0.0, "bertscore_recall": 0.0, "bertscore_f1": 0.0})
    else:
        metrics.update({"bertscore_precision": 0.0, "bertscore_recall": 0.0, "bertscore_f1": 0.0})

    if ROUGE_AVAILABLE:
        try:
            rouge = get_rouge_scorer().score(ground_truth, generated)
            metrics.update({"rouge1_f1": rouge['rouge1'].fmeasure, "rouge2_f1": rouge['rouge2'].fmeasure,
                             "rougeL_f1": rouge['rougeL'].fmeasure})
        except Exception as e:
            logger.warning(f"ROUGE error: {e}")
            metrics.update({"rouge1_f1": 0.0, "rouge2_f1": 0.0, "rougeL_f1": 0.0})
    else:
        metrics.update({"rouge1_f1": 0.0, "rouge2_f1": 0.0, "rougeL_f1": 0.0})

    # BLEU
    if NLTK_METRICS_AVAILABLE and generated.strip() and ground_truth.strip():
        try:
            metrics["bleu_score"] = float(sentence_bleu(
                [ground_truth.split()], generated.split(), smoothing_function=SmoothingFunction().method4))
        except: metrics["bleu_score"] = 0.0
    else:
        metrics["bleu_score"] = 0.0

    # METEOR
    if NLTK_METRICS_AVAILABLE and generated.strip() and ground_truth.strip():
        try: metrics["meteor_score"] = float(meteor_score([ground_truth.split()], generated.split()))
        except: metrics["meteor_score"] = 0.0
    else:
        metrics["meteor_score"] = 0.0

    return metrics


def calculate_semantic_similarity(a: str, b: str) -> float:
    if not a.strip() or not b.strip(): return 0.0
    try:
        e1 = np.array(_embed([a])[0]).reshape(1,-1)
        e2 = np.array(_embed([b])[0]).reshape(1,-1)
        return float(cosine_similarity(e1, e2)[0][0])
    except: return 0.0


# ==================================================
# RAG SYSTEM
# ==================================================

class RAGSystem:
    def __init__(self, qdrant_url: str, collection: str):
        self.retriever = HybridRetriever(
            qdrant_url=qdrant_url, collection_name=collection,
            use_ollama=True, use_reranker=True, embedder=_embedder,
        )

    def retrieve(self, query: str, top_k: int = TOP_K):
        results = self.retriever.search(query, top_k=top_k)
        # for i, r in enumerate(results[:10], 1):
        #     logger.info(
        #         f"TOP{i}: id={r.id} filename={r.metadata.get('filename')} "
        #         f"dense={getattr(r, 'dense_score', None)} sparse={getattr(r,'sparse_score', None)} "
        #         f"rerank={getattr(r,'rerank_score', None)} fused={r.score:.6f}"
        #     )
        docs = [{"content": r.content, "score": r.score,
                 "source": r.metadata.get("filename","unknown"),
                 "chunk":  r.metadata.get("chunk_idx", 0)} for r in results]
        sources = [d["source"] for d in docs]
        parts   = [f"[Source {i}: {d['source']} | chunk {d['chunk']} | score {d['score']:.4f}]\n{d['content']}"
                   for i, d in enumerate(docs, 1)]
        context = "\n\n---\n\n".join(parts)
        if len(context) > MAX_CONTEXT_CHARS:       # BUG 7 FIX
            context = context[:MAX_CONTEXT_CHARS] + "\n\n[Context truncated]"
        return docs, context, sources

    def build_prompt(self, question: str, context: str) -> str:
        return (
            "You are a precise, knowledgeable assistant. Your job is to answer the question "
            "using ONLY the information in the context below.\n\n"
            "RULES:\n"
            "1. Read every chunk carefully before answering.\n"
            "2. If the answer spans multiple chunks, combine them.\n"
            "3. Be specific — include numbers, names, dates, and steps when present.\n"
            "4. If the context genuinely does not contain the answer, say exactly: "
            "'The provided documents do not contain enough information to answer this question.'\n"
            "5. Do NOT say 'not found' if the answer is partially present — give what you have.\n"
            "6. Answer in the same language as the question.\n\n"
            f"CONTEXT:\n{context}\n\n"
            f"QUESTION: {question}\n\n"
            "ANSWER (be thorough and complete):"
        )


# ==================================================
# GROUND TRUTH MATCHING
# ==================================================

def find_ground_truth(question: str, gt_data: dict) -> Optional[dict]:
    """Exact → case-insensitive → fuzzy (threshold 0.75)."""
    if not gt_data or not gt_data.get("questions"):
        return None
    qs = gt_data["questions"]
    qn = question.strip().lower()

    for e in qs:
        if e.get("question","").strip() == question.strip():
            return e
    for e in qs:
        if e.get("question","").strip().lower() == qn:
            return e

    best, best_e = 0.0, None
    for e in qs:
        r = difflib.SequenceMatcher(None, qn, e.get("question","").lower()).ratio()
        if r > best: best, best_e = r, e

    if best >= 0.75:
        logger.info(f"GT fuzzy match ratio={best:.3f}")
        return best_e
    return None


# ==================================================
# SHARED EVALUATION HELPER
# ==================================================

def _build_eval_result(
    model_key: str, answer: str, latency: float,
    question: str, context: str, reference: str,
    docs: list, retrieved_sources: list,
    judge_model: str = JUDGE_GPT,
) -> dict:
    """Evaluate answer against reference using judge + BERTScore + semantic sim."""
    def _judge():
        return evaluate_with_judge(judge_model, question=question, context=context,
                                   answer=answer, ground_truth=reference)
    def _fast():
        return calculate_fast_metrics(answer, reference) if reference.strip() else {}
    def _sem():
        return calculate_semantic_similarity(answer, reference) if reference.strip() else 0.0

    with ThreadPoolExecutor(max_workers=3) as ex:
        fj = ex.submit(_judge); ff = ex.submit(_fast); fs = ex.submit(_sem)
        judge = fj.result(); fast = ff.result(); sem = fs.result()

    overall = judge.get("faithfulness",0)*0.35 + judge.get("completeness",0)*0.35 + judge.get("correctness",0)*0.30
    uniq    = list(dict.fromkeys(retrieved_sources))

    return {
        "model":               model_key,
        "answer":              answer,
        "latency_sec":         round(latency, 2),
        "faithfulness":        judge.get("faithfulness", 0),
        "completeness":        judge.get("completeness", 0),
        "correctness":         judge.get("correctness",  0),
        "overall_score":       round(overall, 2),
        "judge_explanation":   judge.get("explanation",""),
        "semantic_similarity": round(sem, 4),
        "bertscore_f1":        round(fast.get("bertscore_f1",        0.0), 4),
        "bertscore_precision": round(fast.get("bertscore_precision",  0.0), 4),
        "bertscore_recall":    round(fast.get("bertscore_recall",     0.0), 4),
        "rouge1_f1":           round(fast.get("rouge1_f1",  0.0), 4),
        "rouge2_f1":           round(fast.get("rouge2_f1",  0.0), 4),
        "rougeL_f1":           round(fast.get("rougeL_f1",  0.0), 4),
        "bleu_score":          round(fast.get("bleu_score", 0.0), 4),
        "meteor_score":        round(fast.get("meteor_score",0.0), 4),
        "chunks_retrieved":    len(retrieved_sources),
        "unique_files":        len(uniq),
        "avg_chunk_score":     round(sum(d["score"] for d in docs)/len(docs), 4) if docs else 0.0,
        "top_file":            uniq[0] if uniq else "none",
    }


# ==================================================
# SINGLE-MODEL COMPARE (existing app flow)
# ==================================================

def run_rag_evaluation(
    selected_model: str, user_question: str, collection_name: str,
    docs_path, gt_entry: Optional[dict] = None,
    judge_model: str = JUDGE_GPT,
    preindexed: bool = False,           # ← NEW: skip create/index if True
) -> Dict[str, Any]:
    if not preindexed:
        client      = create_temp_collection(collection_name)
        chunk_count = index_uploaded_files(client, collection_name, docs_path)
        if chunk_count == 0:
            raise ValueError("No content extracted. Supported: PDF, DOCX, TXT, MD, CSV, XLSX, JSON, HTML, PPTX, XML, YAML.")

    rag  = RAGSystem(QDRANT_URL, collection_name)
    docs, context, sources = rag.retrieve(user_question, top_k=TOP_K)
    if not docs:
        raise ValueError("Retrieval returned no results.")

    prompt = rag.build_prompt(user_question, context)

    model_answer = ""; gemini_answer = ""; ml = 0.0; gl = 0.0

    def _om():
        nonlocal model_answer, ml
        t0 = time.time(); model_answer  = generate_with_ollama(selected_model, prompt); ml = time.time()-t0
    def _gm():
        nonlocal gemini_answer, gl
        t0 = time.time(); gemini_answer = generate_with_gemini(prompt);                gl = time.time()-t0

    with ThreadPoolExecutor(max_workers=2) as ex:
        f1 = ex.submit(_om); f2 = ex.submit(_gm); f1.result(); f2.result()

    gt_ans      = gt_entry.get("answer","") if gt_entry else ""
    reference   = gt_ans if gt_ans else gemini_answer
    ref_label   = "Ground Truth" if gt_ans else "Gemini Answer"
    eval_answer = model_answer if model_answer else gemini_answer

    res = _build_eval_result(selected_model, eval_answer, ml, user_question, context, reference, docs, sources,
                              judge_model=judge_model)

    # Semantic similarity vs Gemini (always)
    sem_vs_gemini = calculate_semantic_similarity(eval_answer, gemini_answer) if model_answer else 1.0

    uniq = list(dict.fromkeys(sources))
    # delete_temp_collection(collection_name)
    if not preindexed:
        delete_temp_collection(collection_name)

    metrics = {**res,
               "semantic_similarity_vs_gemini": round(sem_vs_gemini, 4),
               "reference_label":  ref_label,
               "model_latency_sec": round(ml, 2),
               "gemini_latency_sec": round(gl, 2),
               "has_ground_truth": bool(gt_ans)}
    metrics.pop("answer", None)

    return {"model_answer": model_answer, "gemini_answer": gemini_answer,
            "model_sources": sources, "metrics": metrics, "gt_entry": gt_entry}


# ==================================================
# MULTI-MODEL COMPARISON
# ==================================================

def run_multi_model_comparison(
    model_names: List[str], user_question: str, docs_path,
    gt_entry: Optional[dict] = None,
    judge_model: str = JUDGE_GPT,
    existing_collection: Optional[str] = None,   # ← NEW: reuse pre-indexed collection
) -> Dict[str, Any]:
    """..."""
    gt_answer = gt_entry.get("answer","") if gt_entry else ""
    has_gt    = bool(gt_answer.strip())

    # Index — reuse existing collection if provided
    if existing_collection:
        cname = existing_collection
    else:
        cname  = f"temp_multi_{_uuid_mod.uuid4().hex[:8]}"
        client = create_temp_collection(cname)
        n      = index_uploaded_files(client, cname, docs_path)
        if n == 0:
            delete_temp_collection(cname)
            raise ValueError("No content extracted from uploaded files.")

    # Retrieve
    rag = RAGSystem(QDRANT_URL, cname)
    docs, context, sources = rag.retrieve(user_question, top_k=TOP_K)
    if not docs:
        delete_temp_collection(cname)
        raise ValueError("Retrieval returned no results.")

    prompt = rag.build_prompt(user_question, context)

    # Generate — Gemini async, Ollama sequential
    generation: Dict[str, Dict] = {}

    def _gen_gemini():
        t0 = time.time(); a = generate_with_gemini(prompt)
        return {"answer": a, "latency": time.time()-t0}

    with ThreadPoolExecutor(max_workers=1) as gex:
        gf = gex.submit(_gen_gemini)
        for name in model_names:
            logger.info(f"Generating: {name}")
            try:
                t0 = time.time(); ans = generate_with_ollama(name, prompt)
                generation[name] = {"answer": ans, "latency": time.time()-t0}
            except Exception as e:
                generation[name] = {"answer": f"Error: {e}", "latency": 0.0}
        try:
            generation["Gemini"] = gf.result()
        except Exception as e:
            generation["Gemini"] = {"answer": f"Error: {e}", "latency": 0.0}

    gemini_answer = generation["Gemini"]["answer"]
    ref_ollama    = gt_answer if has_gt else gemini_answer   # for Ollama models
    ref_gemini    = gt_answer                                  # for Gemini (empty = context-only judge)

    # Evaluate all in parallel
    eval_results: Dict[str, Dict] = {}

    def _eval(key: str):
        gen = generation.get(key, {})
        ref = ref_gemini if key == "Gemini" else ref_ollama
        return _build_eval_result(key, gen.get("answer",""), gen.get("latency",0.0),
                                   user_question, context, ref, docs, sources,
                                   judge_model=judge_model)

    all_keys = list(model_names) + ["Gemini"]
    # with ThreadPoolExecutor(max_workers=min(len(all_keys), 3)) as ex:
    #     futures = {ex.submit(_eval, k): k for k in all_keys}
    #     for f in as_completed(futures):
    #         k = futures[f]
    #         try:    eval_results[k] = f.result()
    #         except Exception as e:
    #             logger.error(f"Eval failed {k}: {e}")
    #             eval_results[k] = {"model": k, "overall_score": 0, "error": str(e)}

    # delete_temp_collection(cname)

    with ThreadPoolExecutor(max_workers=min(len(all_keys), 3)) as ex:
        futures = {ex.submit(_eval, k): k for k in all_keys}
        for f in as_completed(futures):
            k = futures[f]
            try:    eval_results[k] = f.result()
            except Exception as e:
                logger.error(f"Eval failed {k}: {e}")
                eval_results[k] = {"model": k, "overall_score": 0, "error": str(e)}

    # ── Compute both semantic similarities separately ──────────────────────────
    # semantic_similarity_vs_gemini: every model vs Gemini's answer (always)
    # semantic_similarity_vs_gt:     every model vs GT answer (only when GT present)
    gemini_gen_ans = generation.get("Gemini", {}).get("answer", "")
    for key in all_keys:
        gen_ans = generation.get(key, {}).get("answer", "")
        # vs Gemini — Gemini compared with itself = 1.0
        if key == "Gemini":
            eval_results[key]["semantic_similarity_vs_gemini"] = 1.0
        else:
            eval_results[key]["semantic_similarity_vs_gemini"] = round(
                calculate_semantic_similarity(gen_ans, gemini_gen_ans), 4
            )
        # vs GT — only populated when GT is available
        if has_gt and gt_answer.strip():
            eval_results[key]["semantic_similarity_vs_gt"] = round(
                calculate_semantic_similarity(gen_ans, gt_answer), 4
            )
        else:
            eval_results[key]["semantic_similarity_vs_gt"] = None
    # ──────────────────────────────────────────────────────────────────────────

    # delete_temp_collection(cname)
    if not existing_collection:
        delete_temp_collection(cname)

    return {
        "results":           eval_results,
        "gemini_answer":     gemini_answer,
        "ground_truth":      gt_entry,
        "has_ground_truth":  has_gt,
        "reference_label":   "Ground Truth" if has_gt else "Gemini Answer",
        "retrieved_sources": sources,
    }


# ==================================================
# DOCUMENT PREVIEW
# ==================================================

SUPPORTED_EXTENSIONS = {
    '.txt','md','.md','.yaml','.yml','.pdf','.docx','.doc',
    '.csv','.xlsx','.xls','.json','.html','.htm','.pptx','.ppt','.xml',
}

def get_document_preview(docs_path) -> List[Dict]:
    """
    Lightweight scan of uploaded files — returns list of file info dicts.
    Does NOT read file content; only inspects metadata.
    """
    results = []
    for file in Path(docs_path).rglob("*"):
        if not file.is_file(): continue
        if any(p.startswith(".") or p == "__MACOSX" for p in file.parts): continue
        if file.suffix.lower() not in SUPPORTED_EXTENSIONS and file.suffix.lower() != "": continue
        if file.suffix.lower() == ".zip": continue
        size_bytes = file.stat().st_size
        size_str   = (f"{size_bytes/1024:.1f} KB" if size_bytes < 1_048_576
                      else f"{size_bytes/1_048_576:.1f} MB")
        results.append({
            "File":      file.name,
            "Type":      file.suffix.upper().lstrip(".") or "—",
            "Size":      size_str,
            "Size_bytes": size_bytes,
        })
    return sorted(results, key=lambda x: x["Size_bytes"], reverse=True)


# ==================================================
# EXCEL EXPORT
# ==================================================

def export_to_excel(
    question:      str,
    model_answer:  str          = "",
    gemini_answer: str          = "",
    metrics:       dict         = None,
    gt_entry:      dict         = None,
    comparison:    dict         = None,   # full multi-model comparison_data dict
    judge_model:   str          = JUDGE_GPT,
) -> bytes:
    """
    Build an Excel workbook in memory and return its bytes.
    Works for both single-model and multi-model results.
    """
    import io
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()

    # ── Colour helpers ────────────────────────────────────────────────────────
    def _fill(hex_color: str) -> PatternFill:
        return PatternFill("solid", fgColor=hex_color)

    HDR_FILL  = _fill("2563EB")   # blue header
    ALT_FILL  = _fill("EFF6FF")   # light blue alternate row
    GT_FILL   = _fill("DCFCE7")   # green for GT
    GOOD_FILL = _fill("BBF7D0")   # green for better-than-gemini
    BAD_FILL  = _fill("FECACA")   # red for worse
    SIM_FILL  = _fill("FEF9C3")   # yellow for similar
    GEM_FILL  = _fill("DBEAFE")   # blue for gemini baseline

    HDR_FONT  = Font(bold=True, color="FFFFFF", name="Calibri")
    BOLD      = Font(bold=True, name="Calibri")
    NORMAL    = Font(name="Calibri")

    def _hdr_row(ws, row_data: list, row: int, fill=HDR_FILL, font=HDR_FONT):
        for ci, val in enumerate(row_data, 1):
            cell = ws.cell(row=row, column=ci, value=val)
            cell.fill = fill; cell.font = font
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)

    def _auto_width(ws, min_w=12, max_w=50):
        for col in ws.columns:
            length = max((len(str(c.value or "")) for c in col), default=min_w)
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max(length+2, min_w), max_w)

    # ──────────────────────────────────────────────────────────────────────────
    # Sheet 1: Summary
    # ──────────────────────────────────────────────────────────────────────────
    ws1 = wb.active; ws1.title = "Summary"
    ws1.row_dimensions[1].height = 30

    _hdr_row(ws1, ["RAG Evaluation Report"], 1)
    ws1.merge_cells("A1:B1")

    rows_s = [
        ("Question",    question),
        ("Judge Model", judge_model),
        ("GT Available", "Yes" if gt_entry else "No"),
    ]
    if gt_entry:
        rows_s.append(("GT Answer", gt_entry.get("answer","")[:500]))
        rows_s.append(("Source Document", gt_entry.get("source_document","")))
        rows_s.append(("Difficulty", gt_entry.get("difficulty","")))
        rows_s.append(("Question Type", gt_entry.get("question_type","")))

    for ri, (k, v) in enumerate(rows_s, 2):
        ws1.cell(ri, 1, k).font = BOLD
        ws1.cell(ri, 2, str(v)).font = NORMAL
        ws1.cell(ri, 2).alignment = Alignment(wrap_text=True)

    _auto_width(ws1)

    # ──────────────────────────────────────────────────────────────────────────
    # Sheet 2: Answers
    # ──────────────────────────────────────────────────────────────────────────
    ws2 = wb.create_sheet("Answers")
    _hdr_row(ws2, ["Model", "Answer"], 1)
    r = 2
    if comparison:
        for key, res in comparison.get("results", {}).items():
            ws2.cell(r, 1, key).font = BOLD
            ws2.cell(r, 2, res.get("answer","")).alignment = Alignment(wrap_text=True)
            r += 1
    else:
        if model_answer:
            ws2.cell(r, 1, "Model").font = BOLD
            ws2.cell(r, 2, model_answer).alignment = Alignment(wrap_text=True)
            r += 1
        if gemini_answer:
            ws2.cell(r, 1, "Gemini").font = BOLD
            ws2.cell(r, 2, gemini_answer).alignment = Alignment(wrap_text=True)
            r += 1
        if gt_entry:
            ws2.cell(r, 1, "Ground Truth").font = BOLD
            c = ws2.cell(r, 2, gt_entry.get("answer",""))
            c.fill = GT_FILL; c.alignment = Alignment(wrap_text=True)

    ws2.column_dimensions["A"].width = 20
    ws2.column_dimensions["B"].width = 80

    # ──────────────────────────────────────────────────────────────────────────
    # Sheet 3: Metrics
    # ──────────────────────────────────────────────────────────────────────────
    ws3 = wb.create_sheet("Metrics")

    if comparison:
        # Multi-model comparison table
        all_results = comparison.get("results", {})
        has_gt_     = comparison.get("has_ground_truth", False)
        GEMINI      = "Gemini"
        gem_data    = all_results.get(GEMINI, {})
        ordered     = [GEMINI] + sorted(k for k in all_results if k != GEMINI)

        metric_keys = [
            ("overall_score","Overall Score"),("faithfulness","Faithfulness"),
            ("completeness","Completeness"),("correctness","Correctness"),
            ("semantic_similarity","Semantic Sim."),("bertscore_f1","BERTScore F1"),
        ]
        if has_gt_:
            metric_keys += [("rougeL_f1","ROUGE-L"),("bleu_score","BLEU"),("meteor_score","METEOR")]
        metric_keys += [("latency_sec","Latency (s)"),("chunks_retrieved","Chunks"),("avg_chunk_score","Avg Score")]

        hdr = ["Model"] + [label for _, label in metric_keys]
        _hdr_row(ws3, hdr, 1)

        THRESH = {"overall_score":0.4,"faithfulness":0.4,"completeness":0.4,"correctness":0.4,
                  "semantic_similarity":0.02,"bertscore_f1":0.02,"rougeL_f1":0.02,
                  "bleu_score":0.02,"meteor_score":0.02,"latency_sec":1.0}
        HIB    = {k: True for k in THRESH}; HIB["latency_sec"] = False

        for ri, key in enumerate(ordered, 2):
            res = all_results.get(key, {})
            ws3.cell(ri, 1, key).font = BOLD
            if key == GEMINI:
                ws3.row_dimensions[ri].fill = GEM_FILL
                ws3.cell(ri, 1).fill = GEM_FILL
            for ci, (mkey, _) in enumerate(metric_keys, 2):
                val = res.get(mkey, 0)
                cell = ws3.cell(ri, ci, round(val,4) if isinstance(val,float) else val)
                cell.font = NORMAL; cell.alignment = Alignment(horizontal="center")
                if key != GEMINI:
                    g_val = gem_data.get(mkey, 0)
                    try:
                        diff = (val-g_val) if HIB.get(mkey,True) else (g_val-val)
                        thr  = THRESH.get(mkey, 0.5)
                        if diff > thr:    cell.fill = GOOD_FILL
                        elif diff < -thr: cell.fill = BAD_FILL
                        else:             cell.fill = SIM_FILL
                    except: pass
                else:
                    cell.fill = GEM_FILL
    else:
        # Single-model metrics
        if metrics:
            _hdr_row(ws3, ["Metric", "Value"], 1)
            skip = {"judge_explanation","reference_label","has_ground_truth","model","answer","top_file"}
            for ri, (k, v) in enumerate(((k,v) for k,v in metrics.items() if k not in skip), 2):
                ws3.cell(ri, 1, k).font = BOLD
                ws3.cell(ri, 2, round(v,4) if isinstance(v,float) else v)
                if ri % 2 == 0:
                    ws3.cell(ri, 1).fill = ALT_FILL
                    ws3.cell(ri, 2).fill = ALT_FILL

    _auto_width(ws3)

    # ──────────────────────────────────────────────────────────────────────────
    # Serialize
    # ──────────────────────────────────────────────────────────────────────────
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.read()
